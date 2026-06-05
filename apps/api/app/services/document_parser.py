"""Document parsing service for PDF, PPT, DOCX, TXT formats."""

import os
from typing import Dict, List, Optional


class DocumentParser:
    """Parse uploaded documents into plain text content."""

    @staticmethod
    async def parse(file_path: str, content_type: str) -> str:
        """Parse a document file and return its text content.

        Dispatches to the appropriate parser based on content type.
        """
        if not os.path.exists(file_path):
            return ""

        parsers: Dict[str, callable] = {
            "application/pdf": DocumentParser._parse_pdf,
            "application/vnd.openxmlformats-officedocument.presentationml.presentation": DocumentParser._parse_pptx,
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document": DocumentParser._parse_docx,
            "text/plain": DocumentParser._parse_text,
            "text/markdown": DocumentParser._parse_text,
        }

        parser = parsers.get(content_type, DocumentParser._parse_text)
        return await parser(file_path)

    @staticmethod
    async def extract_chunks(
        text: str,
        chunk_size: int = 1000,
        overlap: int = 200,
    ) -> List[dict]:
        """Split text into overlapping chunks for embedding.

        Returns a list of dicts with keys: content, token_count, chunk_index.
        """
        if not text:
            return []

        chunks = []
        start = 0
        index = 0

        while start < len(text):
            end = start + chunk_size
            chunk_text = text[start:end]

            # Try to break at sentence boundary
            if end < len(text):
                last_period = chunk_text.rfind(". ")
                if last_period > chunk_size // 2:
                    chunk_text = chunk_text[: last_period + 1]
                    end = start + last_period + 1

            chunks.append({
                "content": chunk_text.strip(),
                "token_count": len(chunk_text.split()),
                "chunk_index": index,
            })

            start = end - overlap
            index += 1

        return chunks

    @staticmethod
    async def _parse_pdf(file_path: str) -> str:
        """Extract text from a PDF file."""
        try:
            import fitz  # PyMuPDF
        except ImportError:
            # Fallback: return filename-based placeholder
            return f"[PDF content from {os.path.basename(file_path)}]"

        text_parts = []
        doc = fitz.open(file_path)
        for page in doc:
            text_parts.append(page.get_text())
        doc.close()
        return "\n\n".join(text_parts)

    @staticmethod
    async def _parse_pptx(file_path: str) -> str:
        """Extract text from a PowerPoint file."""
        try:
            from pptx import Presentation
        except ImportError:
            return f"[PPTX content from {os.path.basename(file_path)}]"

        text_parts = []
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text_parts.append(paragraph.text)
        return "\n".join(text_parts)

    @staticmethod
    async def _parse_docx(file_path: str) -> str:
        """Extract text from a Word document."""
        try:
            from docx import Document
        except ImportError:
            return f"[DOCX content from {os.path.basename(file_path)}]"

        doc = Document(file_path)
        paragraphs = [para.text for para in doc.paragraphs]
        return "\n\n".join(paragraphs)

    @staticmethod
    async def _parse_text(file_path: str) -> str:
        """Read a plain text or markdown file."""
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            return f.read()


def get_document_parser() -> DocumentParser:
    """Factory function — returns a DocumentParser instance."""
    return DocumentParser()

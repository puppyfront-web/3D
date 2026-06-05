"""File parsers for different document formats.

Extends the base DocumentParser from services with format-specific
parsing logic and metadata extraction.
"""

import os
from typing import Dict, List, Optional

from app.services.document_parser import DocumentParser as BaseParser


class ParsedDocument:
    """Result of parsing a document."""

    def __init__(
        self,
        text: str,
        metadata: Dict[str, any],
        pages: Optional[List[str]] = None,
    ):
        self.text = text
        self.metadata = metadata
        self.pages = pages or [text]

    @property
    def page_count(self) -> int:
        return len(self.pages)

    @property
    def word_count(self) -> int:
        return len(self.text.split())

    @property
    def char_count(self) -> int:
        return len(self.text)


async def parse_file(file_path: str, content_type: str) -> ParsedDocument:
    """Parse a file and return structured content with metadata.

    This is the main entry point for document parsing. It dispatches
    to format-specific parsers and returns rich metadata.
    """
    filename = os.path.basename(file_path)
    file_size = os.path.getsize(file_path) if os.path.exists(file_path) else 0

    base_metadata = {
        "filename": filename,
        "content_type": content_type,
        "file_size": file_size,
        "file_extension": os.path.splitext(filename)[1].lower(),
    }

    parser = BaseParser()
    text = await parser.parse(file_path, content_type)

    # Extract page-level content if possible
    pages = await _extract_pages(file_path, content_type)

    base_metadata.update({
        "word_count": len(text.split()) if text else 0,
        "char_count": len(text) if text else 0,
        "page_count": len(pages),
    })

    return ParsedDocument(text=text, metadata=base_metadata, pages=pages)


async def _extract_pages(file_path: str, content_type: str) -> List[str]:
    """Extract content split by pages."""
    if content_type == "application/pdf":
        return await _extract_pdf_pages(file_path)
    elif content_type in (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ):
        return await _extract_pptx_pages(file_path)
    else:
        # For non-paginated formats, return single "page"
        return []


async def _extract_pdf_pages(file_path: str) -> List[str]:
    """Extract text from PDF pages."""
    try:
        import fitz
    except ImportError:
        return []

    pages = []
    doc = fitz.open(file_path)
    for page in doc:
        pages.append(page.get_text())
    doc.close()
    return pages


async def _extract_pptx_pages(file_path: str) -> List[str]:
    """Extract text from PPTX slides (treated as pages)."""
    try:
        from pptx import Presentation
    except ImportError:
        return []

    pages = []
    prs = Presentation(file_path)
    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    texts.append(para.text)
        pages.append("\n".join(texts))
    return pages

"""PPTX export functionality — converts markdown-like content to PowerPoint slides."""

import os
from datetime import datetime, timezone
from typing import Optional

from app.core.config import settings

# Color constants matching the project theme
TITLE_COLOR = "1E3A5F"  # Deep navy blue
HEADING_COLOR = "2D5A8E"  # Medium blue
BODY_COLOR = "333333"  # Dark gray
ACCENT_COLOR = "00D4FF"  # Electric cyan
LIGHT_GRAY = "999999"


class PPTXExporter:
    """Export content to PowerPoint (.pptx) format."""

    async def export(
        self,
        content: str,
        filename: str,
        title: Optional[str] = None,
        author: str = "3D Wall Platform",
    ) -> str:
        """Export markdown-like content to a PowerPoint presentation.

        Parsing rules:
        - ``# `` → Title slide
        - ``## `` → New slide with heading
        - ``### `` → Sub-heading (same slide)
        - ``- `` / ``* `` → Bullet points
        - ``1. `` ... → Numbered items (rendered as bullets)
        - ``---`` → Slide separator (forces new slide)
        - Other text → Body paragraph

        Returns the file path of the generated presentation.
        """
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        prs = Presentation()
        # Set 16:9 widescreen
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Use the blank layout (no title/subtitle placeholders)
        blank_layout = prs.slide_layouts[6]  # Blank

        lines = content.split("\n")
        current_slide = None
        body_items: list[str] = []

        def _rgb(hex_color: str) -> "RGBColor":
            return RGBColor.from_string(hex_color)

        def _flush_body():
            """Flush accumulated body items onto the current slide."""
            nonlocal body_items
            if not body_items or not current_slide:
                return
            txBox = current_slide.shapes.add_textbox(
                Inches(0.8), Inches(3.0), Inches(11.7), Inches(4.0)
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            for i, item in enumerate(body_items):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.space_after = Pt(6)
                run = p.add_run()
                if item.startswith("- ") or item.startswith("* ") or item.startswith("• "):
                    run.text = "• " + item.lstrip("- *• ")
                    p.level = 0
                else:
                    run.text = item
                run.font.size = Pt(16)
                run.font.color.rgb = _rgb(BODY_COLOR)
            body_items = []

        for line in lines:
            stripped = line.strip()

            if stripped.startswith("# "):
                # Title slide
                _flush_body()
                current_slide = prs.slides.add_slide(blank_layout)
                _add_title_slide(current_slide, stripped[2:], author, _rgb)

            elif stripped.startswith("## "):
                # New section slide
                _flush_body()
                current_slide = prs.slides.add_slide(blank_layout)
                _add_section_heading(current_slide, stripped[3:], _rgb)

            elif stripped.startswith("### "):
                # Sub-heading on current slide
                _flush_body()
                if current_slide:
                    _add_sub_heading(current_slide, stripped[4:], _rgb)

            elif stripped.startswith("---"):
                _flush_body()
                current_slide = None  # Force new slide on next content

            elif stripped.startswith(("- ", "* ")):
                body_items.append(stripped)

            elif stripped and stripped[0].isdigit() and ". " in stripped[:4]:
                # Numbered item → bullet
                body_items.append("• " + stripped.split(". ", 1)[1])

            elif stripped.startswith("| "):
                # Skip table lines in PPTX
                continue

            elif stripped:
                body_items.append(stripped)

        _flush_body()

        # Ensure at least one slide exists
        if len(prs.slides) == 0:
            slide = prs.slides.add_slide(blank_layout)
            _add_title_slide(slide, title or filename, author, _rgb)

        # Save to storage
        storage_dir = os.path.abspath(settings.storage_path)
        os.makedirs(storage_dir, exist_ok=True)
        filepath = os.path.join(storage_dir, filename)
        prs.save(filepath)

        return filepath


def _add_title_slide(slide, title_text: str, author: str, _rgb) -> None:
    """Add a title slide with the project's deep blue theme."""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    # Title
    txBox = slide.shapes.add_textbox(
        Inches(0.8), Inches(2.0), Inches(11.7), Inches(2.5)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = _rgb(TITLE_COLOR)

    # Subtitle / author line
    txBox2 = slide.shapes.add_textbox(
        Inches(0.8), Inches(4.8), Inches(11.7), Inches(1.0)
    )
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.LEFT
    run2 = p2.add_run()
    now_str = datetime.now(timezone.utc).strftime("%Y-%m-%d")
    run2.text = f"{author}  |  {now_str}"
    run2.font.size = Pt(14)
    run2.font.color.rgb = _rgb(LIGHT_GRAY)

    # Accent line
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(0.8), Inches(4.5), Inches(3.0), Inches(0.05),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(ACCENT_COLOR)
    shape.line.fill.background()


def _add_section_heading(slide, heading_text: str, _rgb) -> None:
    """Add a section heading to a slide."""
    from pptx.util import Inches, Pt

    txBox = slide.shapes.add_textbox(
        Inches(0.8), Inches(0.6), Inches(11.7), Inches(1.8)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = heading_text
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = _rgb(HEADING_COLOR)

    # Accent underline
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(0.8), Inches(2.3), Inches(2.0), Inches(0.04),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(ACCENT_COLOR)
    shape.line.fill.background()


def _add_sub_heading(slide, text: str, _rgb) -> None:
    """Add a sub-heading on the current slide."""
    from pptx.util import Inches, Pt

    txBox = slide.shapes.add_textbox(
        Inches(0.8), Inches(2.3), Inches(11.7), Inches(0.6)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = _rgb(TITLE_COLOR)
